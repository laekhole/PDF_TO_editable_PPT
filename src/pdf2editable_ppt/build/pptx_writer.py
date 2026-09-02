"""Package assembly and element emission.

python-pptx owns the parts of a package that are tedious to get right --
content types, relationships, media parts, the theme and master -- and every
shape is appended to the slide as raw DrawingML built by :mod:`drawingml`.

``build_deck`` is re-runnable on purpose: the converter builds a first deck,
verifies it, swaps failed regions for renders, and builds again.
"""

from __future__ import annotations

import io
from typing import Dict, List, Tuple

from lxml import etree
from pptx import Presentation
from pptx.util import Emu

from ..ir import (
    Document,
    Element,
    ElementType,
    ImageContent,
    Page,
    Path,
    TableCell,
    TableContent,
    TextContent,
)
from ..units import PageGeometry, Rect, pt_to_emu, rect_from_points
from . import drawingml as dml

# The slide's own group shape is id 1; everything we add starts above it.
FIRST_SHAPE_ID = 2
BLANK_LAYOUT_INDEX = 6
# A page with nothing to size a deck from gets US Letter.
DEFAULT_PAGE_PT = (612.0, 792.0)

_NAMES = {
    ElementType.TEXT: "TextBox",
    ElementType.IMAGE: "Picture",
    ElementType.LINE: "Line",
    ElementType.RECT: "Rectangle",
    ElementType.ELLIPSE: "Ellipse",
    ElementType.FREEFORM: "Freeform",
    ElementType.TABLE: "Table",
    ElementType.RASTER_FALLBACK: "Fallback",
    ElementType.VECTOR_FALLBACK: "Fallback",
}


class DeckWriter:
    """A deck of one fixed slide size, filled with raw DrawingML fragments."""

    def __init__(self, width_pt: float, height_pt: float) -> None:
        self.prs = Presentation()
        self.prs.slide_width = Emu(pt_to_emu(width_pt))
        self.prs.slide_height = Emu(pt_to_emu(height_pt))
        self._layout = self.prs.slide_layouts[BLANK_LAYOUT_INDEX]
        self._next_id = FIRST_SHAPE_ID

    # ── slides ───────────────────────────────────────────────────────────

    def add_slide(self):
        slide = self.prs.slides.add_slide(self._layout)
        # The Blank layout has no placeholders, but never rely on a template:
        # a slide must hold exactly the source's ink and nothing inherited.
        tree = slide.shapes._spTree
        for sp in list(tree):
            if sp.find(".//{%s}ph" % dml.P_NS) is not None:
                tree.remove(sp)
        return slide

    def set_background(self, slide, color: str) -> None:
        csld = slide._element.find("{%s}cSld" % dml.P_NS)
        existing = csld.find("{%s}bg" % dml.P_NS)
        if existing is not None:
            csld.remove(existing)
        csld.insert(0, self._parse(dml.background_xml(color)))

    def alloc_id(self) -> int:
        value = self._next_id
        self._next_id += 1
        return value

    def append(self, slide, xml: str) -> None:
        """Append one fragment (``p:sp``, ``p:pic``, ``p:graphicFrame``) to the slide."""
        slide.shapes._spTree.append(self._parse(xml))

    def add_image(self, slide, data: bytes) -> str:
        """Embed image bytes as-is and return the relationship id to use."""
        part, r_id = slide.part.get_or_add_image_part(io.BytesIO(data))
        return r_id

    def save(self, path: str) -> None:
        self.prs.save(path)

    @staticmethod
    def _parse(xml: str):
        wrapped = "<w %s>%s</w>" % (dml.NSDECL, xml)
        root = etree.fromstring(wrapped.encode("utf-8"))
        element = root[0]
        root.remove(element)
        etree.cleanup_namespaces(element)
        return element


# ── the document -> deck mapping ───────────────────────────────────────────


def deck_size(document: Document) -> Tuple[float, float]:
    """One deck has one slide size: the largest page, with smaller pages
    letterboxed inside it at 1:1 so no measured size changes."""
    if not document.pages:
        return DEFAULT_PAGE_PT
    return (
        max(p.width_pt for p in document.pages),
        max(p.height_pt for p in document.pages),
    )


def page_geometry(page: Page, deck_w: float, deck_h: float) -> PageGeometry:
    return PageGeometry(
        media_x0=0.0,
        media_y0=0.0,
        media_width=page.width_pt,
        media_height=page.height_pt,
        rotation=0,
        offset_x_pt=(deck_w - page.width_pt) / 2.0,
        offset_y_pt=(deck_h - page.height_pt) / 2.0,
    )


def build_deck(document: Document, output_path: str) -> None:
    deck_w, deck_h = deck_size(document)
    writer = DeckWriter(deck_w, deck_h)
    for page in document.pages:
        slide = writer.add_slide()
        if page.background:
            writer.set_background(slide, page.background)
        geom = page_geometry(page, deck_w, deck_h)
        for el in sorted(page.elements, key=lambda e: e.source_paint_order):
            if el.consumed:
                continue
            emit_element(writer, slide, geom, document, el)
    writer.save(output_path)


def emit_element(writer: DeckWriter, slide, geom: PageGeometry, document: Document, el: Element) -> None:
    kind = el.type
    if kind is ElementType.TEXT:
        _emit_text(writer, slide, geom, el)
    elif kind is ElementType.IMAGE or isinstance(el.content, ImageContent):
        _emit_picture(writer, slide, geom, document, el)
    elif kind is ElementType.LINE:
        _emit_line(writer, slide, geom, el)
    elif kind is ElementType.RECT:
        _emit_rect(writer, slide, geom, el)
    elif kind is ElementType.ELLIPSE:
        _emit_ellipse(writer, slide, geom, el)
    elif kind is ElementType.FREEFORM:
        _emit_freeform(writer, slide, geom, el)
    elif kind is ElementType.TABLE:
        _emit_table(writer, slide, geom, el)
    else:
        # A fallback whose region never got rendered, or a type this writer
        # has no emitter for: say so on the element rather than drop it quietly.
        el.note("nothing was emitted for this %s element" % kind.value)


def _name(el: Element) -> str:
    return "%s %s" % (_NAMES.get(el.type, "Shape"), el.id)


# ── text ─────────────────────────────────────────────────────────────────────


def _emit_text(writer: DeckWriter, slide, geom: PageGeometry, el: Element) -> None:
    content: TextContent = el.content
    if content is None or not content.lines:
        el.note("text element without lines; nothing emitted")
        return
    x, y, cx, cy = geom.rect_to_emu(el.bbox)
    writer.append(
        slide,
        dml.textbox_xml(
            writer.alloc_id(),
            _name(el),
            dml.xfrm(x, y, cx, cy, el.rotation_deg),
            content,
            el.opacity,
        ),
    )


# ── pictures ─────────────────────────────────────────────────────────────────


def _emit_picture(writer: DeckWriter, slide, geom: PageGeometry, document: Document, el: Element) -> None:
    content = el.content
    if not isinstance(content, ImageContent):
        el.note("image element without a rendered asset; nothing emitted")
        return
    asset = document.assets.get(content.asset_id)
    if asset is None or not asset.data:
        el.note("asset %s is missing; nothing emitted" % content.asset_id)
        return
    r_id = writer.add_image(slide, asset.data)
    x, y, cx, cy = geom.rect_to_emu(el.bbox)
    writer.append(
        slide,
        dml.picture_xml(
            writer.alloc_id(),
            _name(el),
            r_id,
            dml.xfrm(x, y, cx, cy, el.rotation_deg, content.flip_h, content.flip_v),
            content.crop,
            el.opacity,
        ),
    )


# ── shapes ───────────────────────────────────────────────────────────────────


def _emit_line(writer: DeckWriter, slide, geom: PageGeometry, el: Element) -> None:
    points = el.content
    try:
        (ax, ay), (bx, by) = points[0], points[1]
    except (TypeError, IndexError, ValueError):
        ax, ay, bx, by = el.bbox.x0, el.bbox.y0, el.bbox.x1, el.bbox.y1
    box = Rect(min(ax, bx), min(ay, by), max(ax, bx), max(ay, by))
    x, y, cx, cy = geom.rect_to_emu(box)
    # The "line" preset runs from the top-left to the bottom-right corner of
    # its box.  Take the endpoint with the smaller x: if it is the *higher*
    # one on the page (larger PDF y) the line already runs that way.
    left_y, right_y = (ay, by) if ax <= bx else (by, ay)
    flip_v = left_y < right_y and abs(ax - bx) > 1e-6
    writer.append(
        slide,
        dml.shape_xml(
            writer.alloc_id(),
            _name(el),
            dml.xfrm(x, y, cx, cy, 0.0, False, flip_v),
            dml.preset_geometry("line"),
            el.style,
            el.opacity,
        ),
    )


def _emit_rect(writer: DeckWriter, slide, geom: PageGeometry, el: Element) -> None:
    spec = el.content if isinstance(el.content, dict) else {}
    prst = spec.get("prst") or "rect"
    x, y, cx, cy = geom.rect_to_emu(el.bbox)
    writer.append(
        slide,
        dml.shape_xml(
            writer.alloc_id(),
            _name(el),
            dml.xfrm(x, y, cx, cy, el.rotation_deg),
            dml.preset_geometry(prst, spec.get("adjust")),
            el.style,
            el.opacity,
        ),
    )


def _emit_ellipse(writer: DeckWriter, slide, geom: PageGeometry, el: Element) -> None:
    x, y, cx, cy = geom.rect_to_emu(el.bbox)
    writer.append(
        slide,
        dml.shape_xml(
            writer.alloc_id(),
            _name(el),
            dml.xfrm(x, y, cx, cy, el.rotation_deg),
            dml.preset_geometry("ellipse"),
            el.style,
            el.opacity,
        ),
    )


def _emit_freeform(writer: DeckWriter, slide, geom: PageGeometry, el: Element) -> None:
    path: Path = el.content
    if not isinstance(path, Path) or not path.segments:
        el.note("freeform without a path; nothing emitted")
        return
    # The frame must contain every point, control points included, so that
    # the local coordinates stay inside the path's own w/h.
    frame = el.bbox
    extent = rect_from_points(path.points())
    if extent is not None:
        frame = frame.union(extent)
    x, y, cx, cy = geom.rect_to_emu(frame)
    writer.append(
        slide,
        dml.shape_xml(
            writer.alloc_id(),
            _name(el),
            dml.xfrm(x, y, cx, cy),
            dml.custom_geometry(path, (x, y), (cx, cy), geom.point_to_emu),
            el.style,
            el.opacity,
        ),
    )


# ── tables ───────────────────────────────────────────────────────────────────


def _emit_table(writer: DeckWriter, slide, geom: PageGeometry, el: Element) -> None:
    table: TableContent = el.content
    if table is None or table.rows < 1 or table.cols < 1:
        el.note("table without a grid; nothing emitted")
        return
    by_pos: Dict[Tuple[int, int], TableCell] = {(c.row, c.col): c for c in table.cells}
    rows_xml: List[str] = []
    for r in range(table.rows):
        cells: List[str] = []
        for c in range(table.cols):
            cell = by_pos.get((r, c))
            if cell is None:
                cell = TableCell(row=r, col=c)
            if cell.merged_by is not None:
                anchor_r, anchor_c = cell.merged_by
                h_merge = c != anchor_c
                v_merge = r != anchor_r
                body = dml.cell_text_body(None)
            else:
                h_merge = v_merge = False
                body = dml.cell_text_body(cell.text, el.opacity)
            cells.append(
                dml.table_cell_xml(
                    body,
                    cell.borders,
                    cell.fill_color,
                    cell.fill_alpha * el.opacity,
                    cell.v_align,
                    cell.margins_pt,
                    h_merge=h_merge,
                    v_merge=v_merge,
                    grid_span=cell.col_span,
                    row_span=cell.row_span,
                )
            )
        rows_xml.append("".join(cells))
    x, y, cx, cy = geom.rect_to_emu(el.bbox)
    widths = table.col_widths_pt or [el.bbox.width / table.cols] * table.cols
    heights = table.row_heights_pt or [el.bbox.height / table.rows] * table.rows
    writer.append(
        slide,
        dml.table_xml(
            writer.alloc_id(),
            _name(el),
            x,
            y,
            cx,
            cy,
            [pt_to_emu(w) for w in widths],
            [pt_to_emu(h) for h in heights],
            rows_xml,
        ),
    )
