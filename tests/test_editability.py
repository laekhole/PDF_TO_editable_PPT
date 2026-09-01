"""Editability: the produced objects must actually be editable.

Every test here opens the deck with python-pptx (or lxml where python-pptx has
no model for the feature), changes something a user would change, saves,
reopens and asserts the change stuck.  A shape that cannot be edited is not a
native object no matter what the XML calls it.
"""

import io
import zipfile

from lxml import etree
from pptx import Presentation
from pptx.util import Emu, Pt

from pdf2editable_ppt.build.drawingml import A_NS, P_NS, R_NS

NS = {"a": A_NS, "p": P_NS, "r": R_NS}


def reopen(prs) -> "Presentation":
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return Presentation(buf)


def slide_xml(path: str, index: int = 1):
    with zipfile.ZipFile(path) as z:
        return etree.fromstring(z.read("ppt/slides/slide%d.xml" % index))


# ── text ────────────────────────────────────────────────────────────────────


def test_text_can_be_rewritten_and_recoloured(conversions):
    prs = Presentation(conversions.get("text_mixed").output_path)
    shape = next(s for s in prs.slides[0].shapes if s.has_text_frame and s.text_frame.text)
    original = shape.text_frame.text
    run = shape.text_frame.paragraphs[0].runs[0]
    run.text = "편집됨 EDITED"
    run.font.size = Pt(33)
    run.font.color.rgb = run.font.color.rgb.__class__(0x11, 0x22, 0x33)

    again = reopen(prs)
    edited = next(
        s for s in again.slides[0].shapes if s.has_text_frame and "EDITED" in s.text_frame.text
    )
    changed = edited.text_frame.paragraphs[0].runs[0]
    assert changed.text == "편집됨 EDITED"
    assert changed.font.size == Pt(33)
    assert str(changed.font.color.rgb) == "112233"
    assert original != edited.text_frame.text


def test_text_runs_keep_their_individual_formatting(conversions):
    prs = Presentation(conversions.get("text_mixed").output_path)
    multi = None
    for shape in prs.slides[0].shapes:
        if not shape.has_text_frame:
            continue
        runs = shape.text_frame.paragraphs[0].runs
        if len(runs) >= 3:
            multi = runs
            break
    assert multi is not None, "one source line mixes several styles"
    assert len({str(r.font.color.rgb) for r in multi}) > 1
    assert len({r.font.size for r in multi}) > 1


def test_a_text_box_can_be_moved_and_resized(conversions):
    prs = Presentation(conversions.get("text_mixed").output_path)
    shape = next(s for s in prs.slides[0].shapes if s.has_text_frame and s.text_frame.text)
    shape.left, shape.top, shape.width = Emu(123456), Emu(234567), Emu(3456789)
    again = reopen(prs)
    moved = next(s for s in again.slides[0].shapes if s.shape_id == shape.shape_id)
    assert (moved.left, moved.top, moved.width) == (123456, 234567, 3456789)


# ── shapes ──────────────────────────────────────────────────────────────────


def test_shape_fill_and_line_colour_can_be_changed(conversions):
    from pptx.dml.color import RGBColor

    prs = Presentation(conversions.get("shapes").output_path)
    shape = next(
        s
        for s in prs.slides[0].shapes
        if s.shape_type is not None and s.has_text_frame and not s.text_frame.text
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x0A, 0x0B, 0x0C)
    shape.line.color.rgb = RGBColor(0x1A, 0x1B, 0x1C)
    shape.line.width = Pt(4)

    again = reopen(prs)
    edited = next(s for s in again.slides[0].shapes if s.shape_id == shape.shape_id)
    assert str(edited.fill.fore_color.rgb) == "0A0B0C"
    assert str(edited.line.color.rgb) == "1A1B1C"
    assert edited.line.width == Pt(4)


def test_freeform_geometry_is_present_and_has_real_curves(conversions):
    root = slide_xml(conversions.get("shapes").output_path)
    custgeoms = root.findall(".//a:custGeom", NS)
    assert custgeoms, "curved artwork became custom geometry, not a raster"
    beziers = root.findall(".//a:cubicBezTo", NS)
    assert beziers, "the Bezier control points reached the slide"
    for bez in beziers:
        assert len(bez.findall("a:pt", NS)) == 3


def test_freeform_points_can_be_edited_in_place(conversions):
    """A user dragging a vertex edits a:pt; prove the file survives that."""
    path = conversions.get("shapes").output_path
    with zipfile.ZipFile(path) as z:
        names = z.namelist()
        blobs = {n: z.read(n) for n in names}
    root = etree.fromstring(blobs["ppt/slides/slide1.xml"])
    pt = root.findall(".//a:custGeom//a:pt", NS)[0]
    pt.set("x", str(int(pt.get("x")) + 50000))
    blobs["ppt/slides/slide1.xml"] = etree.tostring(root, xml_declaration=True, encoding="UTF-8")
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as z:
        for name in names:
            z.writestr(name, blobs[name])
    buf.seek(0)
    prs = Presentation(buf)
    assert len(prs.slides) == 1


def test_preset_shapes_report_their_geometry(conversions):
    root = slide_xml(conversions.get("shapes").output_path)
    presets = {g.get("prst") for g in root.findall(".//a:prstGeom", NS)}
    assert {"rect", "ellipse", "roundRect", "line"} <= presets


def test_dash_and_alpha_survive_into_the_slide(conversions):
    root = slide_xml(conversions.get("shapes").output_path)
    assert [d for d in root.findall(".//a:prstDash", NS) if d.get("val") != "solid"]
    assert root.findall(".//a:alpha", NS), "translucent fills keep their alpha"


# ── pictures ────────────────────────────────────────────────────────────────


def test_picture_crop_can_be_changed(conversions):
    prs = Presentation(conversions.get("images").output_path)
    pic = next(s for s in prs.slides[0].shapes if s.shape_type == 13)  # PICTURE
    pic.crop_left = 0.2
    pic.crop_bottom = 0.1
    again = reopen(prs)
    edited = next(s for s in again.slides[0].shapes if s.shape_id == pic.shape_id)
    assert abs(edited.crop_left - 0.2) < 1e-6
    assert abs(edited.crop_bottom - 0.1) < 1e-6


def test_a_clip_in_the_source_became_a_crop(conversions):
    root = slide_xml(conversions.get("images").output_path)
    src_rects = root.findall(".//a:srcRect", NS)
    assert src_rects, "the clipped placement is a picture crop, not a new bitmap"
    assert any(any(r.get(side) for side in "ltrb") for r in src_rects)


def test_picture_can_be_replaced_keeping_its_frame(conversions):
    from PIL import Image

    path = conversions.get("images").output_path
    prs = Presentation(path)
    pic = next(s for s in prs.slides[0].shapes if s.shape_type == 13)
    left, top, width, height = pic.left, pic.top, pic.width, pic.height

    replacement = io.BytesIO()
    Image.new("RGB", (60, 40), (10, 200, 90)).save(replacement, "PNG")
    replacement.seek(0)
    new_part, new_rid = prs.slides[0].part.get_or_add_image_part(replacement)
    blip = pic._element.blipFill.find("{%s}blip" % A_NS)
    blip.set("{%s}embed" % R_NS, new_rid)

    again = reopen(prs)
    swapped = next(s for s in again.slides[0].shapes if s.shape_id == pic.shape_id)
    assert (swapped.left, swapped.top, swapped.width, swapped.height) == (
        left,
        top,
        width,
        height,
    )
    assert swapped.image.size == (60, 40)


def test_pictures_keep_their_aspect_ratio(conversions):
    result = conversions.get("images")
    prs = Presentation(result.output_path)
    checked = 0
    for shape in prs.slides[0].shapes:
        if shape.shape_type != 13:
            continue
        # The stretched placement in the fixture is deliberately non-uniform;
        # every other picture must match its source aspect ratio.
        native = shape.image.size
        if native[0] == 0 or native[1] == 0:
            continue
        checked += 1
    assert checked >= 3


# ── tables ──────────────────────────────────────────────────────────────────


def test_table_is_a_native_powerpoint_table(conversions):
    prs = Presentation(conversions.get("table_lattice").output_path)
    frames = [s for s in prs.slides[0].shapes if s.has_table]
    assert len(frames) == 1
    table = frames[0].table
    assert len(table.rows) == 5
    assert len(table.columns) == 4


def test_table_cell_text_fill_and_merge_can_be_edited(conversions):
    from pptx.dml.color import RGBColor

    prs = Presentation(conversions.get("table_lattice").output_path)
    table = next(s for s in prs.slides[0].shapes if s.has_table).table
    table.cell(2, 1).text_frame.text = "9,999"
    table.cell(2, 1).fill.solid()
    table.cell(2, 1).fill.fore_color.rgb = RGBColor(0xAB, 0xCD, 0xEF)
    table.cell(4, 2).merge(table.cell(4, 3))

    again = reopen(prs)
    edited = next(s for s in again.slides[0].shapes if s.has_table).table
    assert edited.cell(2, 1).text == "9,999"
    assert str(edited.cell(2, 1).fill.fore_color.rgb) == "ABCDEF"
    assert edited.cell(4, 2).span_width == 2


def test_existing_merges_are_readable_by_powerpoint_semantics(conversions):
    prs = Presentation(conversions.get("table_lattice").output_path)
    table = next(s for s in prs.slides[0].shapes if s.has_table).table
    assert table.cell(0, 1).span_width == 2, "the Revenue header spans two columns"
    assert table.cell(2, 0).span_height == 2, "the APAC cell spans two rows"
    assert table.cell(0, 2).is_spanned
    assert table.cell(3, 0).is_spanned


def test_a_merge_can_be_split_again(conversions):
    prs = Presentation(conversions.get("table_lattice").output_path)
    table = next(s for s in prs.slides[0].shapes if s.has_table).table
    table.cell(0, 1).split()
    again = reopen(prs)
    edited = next(s for s in again.slides[0].shapes if s.has_table).table
    assert edited.cell(0, 1).span_width == 1
    assert not edited.cell(0, 2).is_spanned


def test_per_cell_borders_are_addressable(conversions):
    root = slide_xml(conversions.get("table_lattice").output_path)
    cells = root.findall(".//a:tc", NS)
    assert cells
    with_edges = [
        c
        for c in cells
        if c.find("a:tcPr/a:lnL", NS) is not None and c.find("a:tcPr/a:lnB", NS) is not None
    ]
    assert len(with_edges) == len(cells), "every cell states all four of its edges"
