"""Experimental OCR: local engines only, separate output, nothing hidden."""

import json
import os

import pytest

from conftest import ARTIFACT_DIR, fixture_path
from pdf2editable_ppt import ocr
from pdf2editable_ppt.units import Rect


def _word(text, x0, x1, y0=0.0, y1=10.0, conf=90.0):
    return ocr.OcrWord(text=text, bbox=Rect(x0, y0, x1, y1), confidence=conf)


# ── pure functions ──────────────────────────────────────────────────────────


def test_character_error_rate_ignores_whitespace_and_counts_edits():
    assert ocr.character_error_rate("한국도로공사", "한국도로공사") == 0.0
    assert ocr.character_error_rate("한국 도로 공사", "한국도로공사") == 0.0
    assert ocr.character_error_rate("abcdef", "abXdef") == pytest.approx(1 / 6)
    assert ocr.character_error_rate("", "") == 0.0
    assert ocr.character_error_rate("", "x") == 1.0


def test_syllable_boxes_on_a_one_em_grid_join_into_one_word():
    """Tesseract boxes Hangul syllable by syllable; the pitch says it's one word."""
    em = 10.0
    words = [_word(ch, i * em + 1, i * em + 7) for i, ch in enumerate("한국도로공사")]
    lines = ocr.group_lines(words)
    assert len(lines) == 1
    assert lines[0].text == "한국도로공사"


def test_a_third_of_an_em_extra_pitch_is_a_space():
    """Realistic Hangul ink (~0.85 em wide) on a 10 pt grid, plus a 3.5 pt space."""
    em = 10.0
    left = [_word(ch, i * em + 0.5, i * em + 9) for i, ch in enumerate("한국")]
    start = 2 * em + 3.5
    right = [_word(ch, start + i * em + 0.5, start + i * em + 9) for i, ch in enumerate("도로")]
    lines = ocr.group_lines(left + right)
    assert lines[0].text == "한국 도로"


def test_a_narrow_syllable_does_not_split_a_word():
    """"도" has narrow ink; the ink gap after it looks like a space, the pitch does not."""
    em = 10.0
    boxes = [
        _word("한", 1, 8),
        _word("도", em + 3, em + 6),  # narrow ink, wide ink gap
        _word("로", 2 * em + 1, 2 * em + 8),
    ]
    assert ocr.group_lines(boxes)[0].text == "한도로"


def test_engine_grouped_latin_words_use_the_ink_gap():
    words = [_word("Safety", 0, 30), _word("first", 33.5, 55), _word("now", 59, 75)]
    assert ocr.group_lines(words)[0].text == "Safety first now"
    tight = [_word("Sa", 0, 10), _word("fety", 10.5, 30)]
    assert ocr.group_lines(tight)[0].text == "Safety"


def test_rows_are_split_at_column_gutters():
    words = [_word("좌측", 0, 20), _word("우측", 200, 220)]
    lines = ocr.group_lines(words)
    assert [ln.text for ln in lines] == ["좌측", "우측"]


def test_lines_come_back_top_to_bottom():
    words = [_word("아래", 0, 20, y0=0, y1=10), _word("위", 0, 10, y0=50, y1=60)]
    assert [ln.text for ln in ocr.group_lines(words)] == ["위", "아래"]


def test_low_confidence_long_words_are_kept_and_short_noise_is_dropped(monkeypatch):
    class Fake(ocr.OcrEngine):
        name = "fake"

        def available(self):
            return True

        def recognise(self, image, languages):
            return [
                ("BOT(Build-Operate-Transfer)", (10, 10, 300, 40), 8.0),  # long, low
                ("ㄴㄴ", (10, 60, 40, 80), 12.0),  # short, low: noise
                ("|", (10, 100, 12, 140), 95.0),  # no alphanumerics: noise
                ("정상", (10, 150, 60, 180), 92.0),
            ]

    page = ocr.ocr_page(fixture_path("scanned_korean"), 0, 612, 792, Fake(), dpi=72)
    texts = [ln.text for ln in page.lines]
    assert "BOT(Build-Operate-Transfer)" in texts
    assert "정상" in texts
    dropped = {w.text for w in page.dropped}
    assert dropped == {"ㄴㄴ", "|"}
    sidecar = ocr.page_to_dict(page)
    assert len(sidecar["droppedAsNoise"]) == 2, "what was thrown away is still on record"


def test_pick_engine_returns_none_when_nothing_is_installed(monkeypatch):
    monkeypatch.setattr(ocr.TesseractEngine, "available", lambda self: False)
    monkeypatch.setattr(ocr.PaddleEngine, "available", lambda self: False)
    assert ocr.pick_engine("auto") is None


# ── with a real engine ──────────────────────────────────────────────────────


@pytest.fixture(scope="module")
def tesseract():
    engine = ocr.TesseractEngine()
    if not engine.available():
        pytest.skip("tesseract with Korean data is not installed")
    return engine


def test_korean_scan_is_read_within_the_measured_error_budget(tesseract):
    """The number in docs/testing.md is a test, not a claim."""
    truth = open(fixture_path("scanned_korean")[:-4] + ".truth.txt", encoding="utf-8").read()
    page = ocr.ocr_page(fixture_path("scanned_korean"), 0, 612, 792, tesseract)
    cer = ocr.character_error_rate(truth, page.text)
    assert cer < 0.05, "CER %.3f" % cer
    assert len(page.lines) == 10
    assert page.mean_confidence() > 85


def test_word_boundaries_are_close_to_the_truth(tesseract):
    truth = open(fixture_path("scanned_korean")[:-4] + ".truth.txt", encoding="utf-8").read().splitlines()
    page = ocr.ocr_page(fixture_path("scanned_korean"), 0, 612, 792, tesseract)
    got = [ln.text.count(" ") for ln in page.lines]
    want = [t.count(" ") for t in truth]
    assert len(got) == len(want)
    assert sum(abs(a - b) for a, b in zip(got, want)) <= len(want) * 1.5


def test_experimental_ocr_writes_a_separate_deck_and_leaves_the_main_one_alone(tesseract, conversions):
    import hashlib

    from pdf2editable_ppt.converter import run_experimental_ocr
    from pptx import Presentation

    result = conversions.get("scanned_korean")
    before = hashlib.sha256(open(result.output_path, "rb").read()).hexdigest()
    deck = os.path.join(ARTIFACT_DIR, "scanned_korean.ocr.pptx")
    sidecar = os.path.join(ARTIFACT_DIR, "scanned_korean.ocr.json")
    summary = run_experimental_ocr(result, fixture_path("scanned_korean"), deck, sidecar)

    after = hashlib.sha256(open(result.output_path, "rb").read()).hexdigest()
    assert before == after, "the fidelity deck must not change"
    assert summary["experimental"] is True
    assert summary["pages"][0]["lines"] == 10
    assert result.report["summary"]["ocr"]["engine"] == "tesseract"

    prs = Presentation(deck)
    assert len(prs.slides) == 1
    slide = prs.slides[0]
    assert not any(s.shape_type == 13 for s in slide.shapes), "no scan bitmap under the text"
    texts = [s.text_frame.text for s in slide.shapes if s.has_text_frame]
    assert any("Safety first" in t for t in texts)
    assert any("experimental" in t for t in texts), "the banner says what this deck is"
    notes = slide.notes_slide.notes_text_frame.text
    assert "OCR (experimental)" in notes and "Safety first" in notes

    data = json.load(open(sidecar, encoding="utf-8"))
    assert data["experimental"] is True
    assert data["pages"][0]["lines"][0]["words"][0]["confidence"] > 0


def test_cli_ocr_flag_produces_the_sidecar_files(tesseract, tmp_path):
    from pdf2editable_ppt.cli import main

    out = tmp_path / "scan.pptx"
    code = main(
        [
            fixture_path("scanned_korean"),
            "-o",
            str(out),
            "--mode",
            "fast",
            "--ocr",
            "experimental",
            "--report",
            str(tmp_path / "scan.report.json"),
            "-q",
        ]
    )
    assert code == 0
    assert (tmp_path / "scan.ocr.pptx").exists()
    assert (tmp_path / "scan.ocr.json").exists()
    report = json.loads((tmp_path / "scan.report.json").read_text(encoding="utf-8"))
    assert report["summary"]["ocr"]["experimental"] is True


def test_ocr_is_off_by_default(tmp_path):
    from pdf2editable_ppt.cli import main

    out = tmp_path / "scan.pptx"
    assert main([fixture_path("scanned_korean"), "-o", str(out), "--mode", "fast", "-q"]) == 0
    assert not (tmp_path / "scan.ocr.pptx").exists()
