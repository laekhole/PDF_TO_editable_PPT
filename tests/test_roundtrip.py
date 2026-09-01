"""PPTX -> PDF -> PPTX round trip.

Building the source as a PowerPoint deck first gives an unambiguous ground
truth: we know exactly what shapes, text and table the PDF was made from, so
the rebuilt deck can be compared against the original deck's own objects and
not only against a picture.
"""

import os

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

from conftest import ARTIFACT_DIR
from pdf2editable_ppt.converter import convert
from pdf2editable_ppt.ir import ElementType
from pdf2editable_ppt.pipeline import ConvertOptions
from pdf2editable_ppt.verify import compare as cmpmod
from pdf2editable_ppt.verify import render as rendermod

SLIDE_W = Emu(9144000)   # 10in
SLIDE_H = Emu(6858000)   # 7.5in


def build_source_deck(path: str) -> dict:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Emu(457200), Emu(304800), Emu(6400800), Emu(600000))
    frame = title.text_frame
    frame.word_wrap = False
    para = frame.paragraphs[0]
    run = para.add_run()
    run.text = "Round trip 라운드 트립"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1B, 0x33, 0x77)

    body = slide.shapes.add_textbox(Emu(457200), Emu(1100000), Emu(4000000), Emu(900000))
    body.text_frame.word_wrap = False
    for index, text in enumerate(("First line", "두 번째 줄", "Third line")):
        paragraph = (
            body.text_frame.paragraphs[0] if index == 0 else body.text_frame.add_paragraph()
        )
        paragraph.alignment = PP_ALIGN.LEFT
        run = paragraph.add_run()
        run.text = text
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0x22, 0x22, 0x22)

    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(5200000), Emu(1100000), Emu(2200000), Emu(1000000)
    )
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(0xF2, 0xC0, 0x33)
    rect.line.color.rgb = RGBColor(0x33, 0x44, 0x99)
    rect.line.width = Pt(2)
    rect.text_frame.text = ""

    oval = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Emu(5200000), Emu(2400000), Emu(1600000), Emu(1100000)
    )
    oval.fill.solid()
    oval.fill.fore_color.rgb = RGBColor(0x55, 0xB0, 0xE0)
    oval.line.color.rgb = RGBColor(0x11, 0x55, 0x88)

    rows, cols = 3, 3
    table_shape = slide.shapes.add_table(
        rows, cols, Emu(457200), Emu(2600000), Emu(4200000), Emu(1400000)
    )
    table = table_shape.table
    values = [
        ["Item", "Q1", "Q2"],
        ["APAC", "1,240", "1,455"],
        ["EMEA", "980", "1,102"],
    ]
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.text = values[r][c]
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(12)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                RGBColor(0x1B, 0x33, 0x77) if r == 0 else RGBColor(0xF2, 0xF4, 0xF9)
            )
            if r == 0:
                cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(
                    0xFF, 0xFF, 0xFF
                )
    prs.save(path)
    return {"rows": rows, "cols": cols, "texts": values}


@pytest.fixture(scope="module")
def round_trip(have_renderer):
    if not have_renderer:
        pytest.skip("no PPTX renderer (LibreOffice) available")
    source_pptx = os.path.join(ARTIFACT_DIR, "roundtrip.source.pptx")
    truth = build_source_deck(source_pptx)
    pdf = rendermod.pptx_to_pdf(source_pptx, ARTIFACT_DIR)
    assert pdf, "the source deck must export to PDF"
    rebuilt = os.path.join(ARTIFACT_DIR, "roundtrip.rebuilt.pptx")
    result = convert(pdf, rebuilt, options=ConvertOptions(verify=True))
    return {"truth": truth, "pdf": pdf, "result": result, "source_pptx": source_pptx}


def test_round_trip_keeps_every_character(round_trip):
    integrity = round_trip["result"].report["summary"]["textIntegrity"]
    assert integrity["identical"], integrity


def test_round_trip_rebuilds_the_table_natively(round_trip):
    tables = [
        el
        for page in round_trip["result"].document.pages
        for el in page.elements
        if el.type is ElementType.TABLE and not el.consumed
    ]
    assert len(tables) == 1
    grid = tables[0].content
    assert (grid.rows, grid.cols) == (3, 3)
    by_pos = {(c.row, c.col): c for c in grid.cells}
    assert by_pos[(1, 0)].text.text == "APAC"
    assert by_pos[(2, 2)].text.text == "1,102"
    assert by_pos[(0, 0)].fill_color is not None


def test_round_trip_rebuilds_the_shapes(round_trip):
    kinds = {
        el.type
        for page in round_trip["result"].document.pages
        for el in page.elements
        if not el.consumed
    }
    assert ElementType.RECT in kinds
    assert ElementType.ELLIPSE in kinds
    assert ElementType.TEXT in kinds


def test_round_trip_does_not_fall_back_to_a_page_image(round_trip):
    from pdf2editable_ppt.ir import Outcome

    live = [
        el
        for page in round_trip["result"].document.pages
        for el in page.elements
        if not el.consumed
    ]
    assert all(el.outcome is not Outcome.PAGE_FALLBACK for el in live)


def test_round_trip_looks_like_the_original_deck(round_trip):
    """Compare the rebuilt deck against the ORIGINAL deck, not just the PDF."""
    original = rendermod.render_pptx_pages(round_trip["source_pptx"], dpi=150)
    rebuilt = rendermod.render_pptx_pages(round_trip["result"].output_path, dpi=150)
    assert original and rebuilt
    from PIL import Image

    canvas = Image.new(
        "RGB", (original[0].width, original[0].height * 2 + 8), "#cccccc"
    )
    canvas.paste(original[0], (0, 0))
    canvas.paste(rebuilt[0].resize(original[0].size), (0, original[0].height + 8))
    canvas.save(os.path.join(ARTIFACT_DIR, "roundtrip.compare.png"))

    comparison = cmpmod.compare_images(
        original[0], rebuilt[0], smooth_px=cmpmod.PAGE_THRESHOLDS.smooth_px
    )
    ok, why = cmpmod.evaluate(comparison, cmpmod.PAGE_THRESHOLDS)
    assert ok, "%s (%s)" % (why, comparison.to_dict())
